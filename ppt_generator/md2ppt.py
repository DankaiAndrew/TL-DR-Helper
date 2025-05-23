import json
import os
from io import BytesIO

import markdown
import requests as req
from ppt_generator.config import PPT_DIR, PPT_MODE_DIR, IMG_DESCRIPTION_DIC_PATH
from PIL import Image
from PIL.ImageQt import rgb
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.text.text import Font
from pptx.util import Cm, Inches, Pt
from ppt_generator.utils import read_json, extract_text_using_regex,get_random_file, dump_json, calculate_edit_distance

#from ppt_generator.content_generator.img_search import get_img
from ppt_generator.markdown_parser import Heading, Out, parse_str
import subprocess
import platform

"""
This module contains the PptGenerator class, which is used to convert a markdown string into a PowerPoint presentation.
    PptGenerator: Converts a markdown string into a PowerPoint presentation.
    MarkdownCategory: Contains constants for markdown categories.
    MD2Slide: Converts a markdown element into a PowerPoint slide.
    MD2TitleSlide: Converts a markdown element into a PowerPoint slide using a theme.
"""


def get_os():
    os_name = platform.system()
    if os_name == "Darwin":
        return "macOS"
    elif os_name == "Windows":
        return "Windows"
    else:
        return os_name  # Linux or other
    
os_name=get_os()
class PptGenerator:
    """
    Converts a markdown string into a PowerPoint presentation.
    """

    prs: Presentation = None
    md_str: str = None
    out: Out = None
    tree: Heading = None
    theme: str = None
    operation_system: str = os_name

    def __init__(
        self, client, img_dic, md_str: str, theme_path: str, save_path: str = PPT_DIR + "test.pptx"
    ) -> None:
        self.theme = theme_path
        theme_param_path = os.path.join(self.theme, "mode.json")
        with open(theme_param_path, encoding="utf-8") as f:
            self.theme_param = json.load(f)

        self.init_pptx(theme_path)
        self.init_markdown(md_str)
        self.client = client
        # generate the title page
        MD2TitleSlide(self.prs, self.theme, self.ppt_main_theme)
        # prepare the image for the main page
        self.img_dic = img_dic
        # generate the slides
        self.traverse_tree(self.tree)
        # save the ppt
        self.prs.save(save_path)
        full_path = os.path.abspath(save_path) # get the full path of the ppt
        self.convert_to_pdf_with_powerpoint(ppt_path=full_path, pdf_path=full_path.replace(".pptx", ".pdf"),os_name=os_name)


    def init_pptx(self, theme_path: str = PPT_MODE_DIR + "1") -> None:
        """
        Initializes the PowerPoint presentation.
        """
        prs = Presentation()
        prs.slide_height = Cm(self.theme_param["slide_size"]["height"])
        prs.slide_width = Cm(self.theme_param["slide_size"]["width"])
        self.theme = theme_path
        self.prs = prs

    def init_markdown(self, md_str: str) -> None:
        """
        Initializes the markdown string.
        """
        self.md_str = md_str
        self.out = parse_str(md_str)
        self.tree = self.out.main
        self.ppt_main_theme = self.tree.text


    def traverse_tree(self, heading: Heading) -> None:
        """
        Traverses the markdown tree and converts each heading into a slide.
        """
        if heading is None:
            return

        if not heading.source or not heading.source.strip():
            # When the heading has no content
            content = ""
            if len(heading.children) > 0:
                content = "\n".join(child.text for child in heading.children)
            content = content.strip()

            if heading.text_source[:2] == "# ":
                # First level heading - create Table of Contents
                MD2Slide(self.prs, self.theme, "Table of Contents", content=content)
            elif heading.text_source[:3] == "## ":
                # Second level heading - create a Theme Slide
                MD2ThemeSlide(self.prs, self.theme, title=heading.text, content=content)
            else:
                # Other heading levels (unlikely)
                MD2Slide(self.prs, self.theme, heading.text, content=content)

        else:
            # When the heading has content
            content_list = heading.source.split("\n\n")  # split content by double line break
            content_max_word_num = 200
            begin_index = 0
            input_text_list = []

            while begin_index < len(content_list):
                input_token_i = 0
                input_text = ""
                for i in range(begin_index, len(content_list)):
                    if len(content_list[i]) + input_token_i < content_max_word_num:
                        input_text += "\n" + content_list[i]
                        input_token_i += len(content_list[i])
                        begin_index = i
                    else:
                        if input_token_i == 0:
                            input_text += content_list[i]
                        break
                begin_index += 1
                input_text_list.append(input_text)

            for content_i, input_text in enumerate(input_text_list):
                img_dict = {}  # Placeholder for future image dictionary matching
                title = heading.text if content_i == 0 else ""  # Only first page shows title
                img_path = ''
                for key, value in self.img_dic.items():
                    if calculate_edit_distance(value, heading.text) <= 2:
                        img_path = key
                MD2Slide(
                    self.prs,
                    self.theme,
                    title=title,
                    img_path=img_path,
                    content=input_text.strip(),
                )

        # Recursively traverse the children
        if len(heading.children) > 0:
            for child in heading.children:
                self.traverse_tree(child)


    def convert_to_pdf_with_powerpoint(self, ppt_path: str, pdf_path: str, os_name: str)->None:
        # AppleScript to open PPT, save as PDF, and close
        if os_name == "macOS":
            applescript = f'''
            tell application "Microsoft PowerPoint"
                try
                    set pptPath to POSIX file "{ppt_path}" as string
                    open pptPath
                    set theDoc to active presentation
                    set pdfPath to POSIX file "{pdf_path}" as string
                    save theDoc in pdfPath as save as PDF
                    close theDoc saving no
                on error errMsg number errNum
                    display dialog "Error: " & errMsg & " Error Number: " & errNum
                end try
            end tell
            '''
            # Run the AppleScript
            subprocess.run(["osascript", "-e", applescript], check=True)
        else:
            try:
                from ppt_generator.windows_ppt2pdf import win_ppt2pdf
                win_ppt2pdf(ppt_path, pdf_path)
            except:
                pass

class MarkdownCategory:
    TITLE = "#"
    CONTENT = "<p>"

    pass


class MD2Slide:
    """
    Converts a markdown element into a PowerPoint slide.
    """

    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "Times New Roman"
    font_title_size: Pt = Pt(26)
    font_content_size: Pt = Pt(14)
    font_title_color: rgb = RGBColor(0, 0, 0)
    font_content_color: rgb = RGBColor(0, 0, 0)
    title_box = (Inches(0.3), Inches(0.3), Cm(24.24), Inches(0.8))
    content_box = (Cm(2.54), Cm(4.12), Cm(20.32), Cm(12.70))

    def __init__(
        self, presentation, theme_path, title, content, *args, img_path='', **kwargs
    ):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme, "mode.json")
        with open(self.theme_param_path, encoding="utf-8") as f:
            self.theme_param = json.load(f)
        print(self.theme)
        page_params = self.theme_param["main_page"]

        self.title_box = (
            Cm(page_params["title_info"]["pos_x"]),
            Cm(page_params["title_info"]["pos_y"]),
            Cm(page_params["title_info"]["width"]),
            Cm(page_params["title_info"]["height"]),
        )
        self.content_box = (
            Cm(page_params["content_info"]["pos_x"]),
            Cm(page_params["content_info"]["pos_y"]),
            Cm(page_params["content_info"]["width"]),
            Cm(page_params["content_info"]["height"]),
        )
        self.font_content_size = Pt(page_params["content_info"]["font_size"])
        self.font_title_size = Pt(page_params["title_info"]["font_size"])
        self.font_name = page_params["title_info"]["font_name"]
        if self.title == "Table of Contents":
            page_params = self.theme_param["catalog_page"]
            self.title_box = (
                Cm(page_params["title_info"]["pos_x"]),
                Cm(page_params["title_info"]["pos_y"]),
                Cm(page_params["title_info"]["width"]),
                Cm(page_params["title_info"]["height"]),
            )
            self.font_title_size = Pt(page_params["title_info"]["font_size"])
            self.font_name = page_params["title_info"]["font_name"]

        self.img_theme = theme_path + "/" + "img"
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()
        if page_params.get("img_info") and img_path:
            self.img_path = img_path
            # img_orginal_h = img_dict["height"]
            # img_orginal_w = img_dict["width"]
            # img_h = float(page_params["img_info"]["width"])*float(img_orginal_h)/float(img_orginal_w)
            self.img_box = (
                Cm(page_params["img_info"]["pos_x"]),
                Cm(page_params["img_info"]["pos_y"]),
                Cm(page_params["img_info"]["width"]),
                Cm(page_params["img_info"]["height"]),
            )
            self.init_img()

    def init_slide(self):
        img_box = (
            Cm(0),
            Cm(0),
            self.presentation.slide_width,
            self.presentation.slide_height,
        )
        # add picture
        img_path = get_random_file(self.img_theme)
        picture = self.slide.shapes.add_picture(img_path, *img_box)
        # picture.left = 0
        # picture.top = 0
        # picture.width = self.presentation.slide_width
        # picture.height = self.presentation.slide_height

    def init_img(self):
        picture = self.slide.shapes.add_picture(self.img_path, *self.img_box)

    def init_font(self, **kwargs):
        if "font_name" in kwargs:
            self.font_name = kwargs["font_name"]
        if "font_title_size" in kwargs:
            self.font_title_size = kwargs["font_title_size"]
        if "font_content_size" in kwargs:
            self.font_content_size = kwargs["font_content_size"]
        if "font_title_color" in kwargs:
            self.font_title_color = kwargs["font_title_color"]
        if "font_content_color" in kwargs:
            self.font_content_color = kwargs["font_content_color"]
        if "content_box" in kwargs:
            self.content_box = kwargs["content_box"]
        if "title_box" in kwargs:
            self.title_box = kwargs["title_box"]

    def get_font(self, font: Font, category: str):
        """
        Sets the font properties for a given category.
        """
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.bold = True
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(*self.title_box)
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # add title
        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(*self.content_box)
        tf = text_box_content.text_frame
        tf.clear()  # Clear existing content
        # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        # main content
        paragraph = tf.paragraphs[0]
        self.content = self.content.replace("<p>", "").replace("</p>", "\n")
        self.content = self.content.replace("\n\n", "\n").replace("\n", "\n\n")
        paragraph.text = self.content
        self.processing_md_str(self.content.replace("<p>", "").replace("</p>", "\n"))
        # paragraph.text = self.content
        self.get_font(paragraph.font, MarkdownCategory.CONTENT)
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def processing_md_str(self, md_str):
        print(md_str)
        md = markdown.Markdown()
        html1 = md.convert(md_str)
        print(html1)

class MD2ThemeSlide:
    """
    Converts a markdown theme section into a PowerPoint theme slide using a theme.
    """

    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "Times New Roman"
    font_title_size: Pt = Pt(40)
    font_content_size: Pt = Pt(20)
    font_title_color: RGBColor = RGBColor(0, 0, 0)
    font_content_color: RGBColor = RGBColor(0, 0, 0)
    title_box = (Cm(0), Cm(0), Cm(0), Cm(0))
    content_box = (Cm(0), Cm(0), Cm(0), Cm(0))

    def __init__(self, presentation, theme_path, title, content="", *args, img_path='', **kwargs):
        self.presentation = presentation
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.title = title
        self.content = content
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme, "mode.json")
        with open(self.theme_param_path, encoding="utf-8") as f:
            self.theme_param = json.load(f)
        print(self.theme)

        theme_page_params = self.theme_param["theme_page"]

        self.title_box = (
            Cm(theme_page_params["title_info"]["pos_x"]),
            Cm(theme_page_params["title_info"]["pos_y"]),
            Cm(theme_page_params["title_info"]["width"]),
            Cm(theme_page_params["title_info"]["height"]),
        )
        self.content_box = (
            Cm(theme_page_params["content_info"]["pos_x"]),
            Cm(theme_page_params["content_info"]["pos_y"]),
            Cm(theme_page_params["content_info"]["width"]),
            Cm(theme_page_params["content_info"]["height"]),
        )
        self.font_title_size = Pt(theme_page_params["title_info"]["font_size"])
        self.font_content_size = Pt(theme_page_params["content_info"]["font_size"])
        self.font_name = theme_page_params["title_info"]["font_name"]
        self.img_theme = theme_path + "/" + "img"

        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()
        self.init_content()

    def init_slide(self):
        img_box = (
            Cm(0),
            Cm(0),
            self.presentation.slide_width,
            self.presentation.slide_height,
        )
        img_path = get_random_file(self.img_theme)
        picture = self.slide.shapes.add_picture(img_path, *img_box)

    def init_font(self, **kwargs):
        if "font_name" in kwargs:
            self.font_name = kwargs["font_name"]
        if "font_title_size" in kwargs:
            self.font_title_size = kwargs["font_title_size"]
        if "font_content_size" in kwargs:
            self.font_content_size = kwargs["font_content_size"]
        if "font_title_color" in kwargs:
            self.font_title_color = kwargs["font_title_color"]
        if "font_content_color" in kwargs:
            self.font_content_color = kwargs["font_content_color"]
        if "content_box" in kwargs:
            self.content_box = kwargs["content_box"]
        if "title_box" in kwargs:
            self.title_box = kwargs["title_box"]

    def get_font(self, font: Font, category: str):
        """
        Sets the font properties for a given category.
        """
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.bold = True
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
        elif category == MarkdownCategory.CONTENT:
            font.size = self.font_content_size
            font.color.rgb = self.font_content_color

    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(*self.title_box)
        tf = text_box.text_frame
        tf.clear()
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        paragraph = tf.paragraphs[0]
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    def init_content(self):
        shapes = self.slide.shapes
        text_box_content = shapes.add_textbox(*self.content_box)
        tf = text_box_content.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

        # split the content into bullet points
        for line in self.content.split("\n"):
            if line.strip() == "":
                continue
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 0  # bullet point
            self.get_font(p.font, MarkdownCategory.CONTENT)


class MD2TitleSlide:
    """
    Converts a markdown element into a PowerPoint slide using a theme.
    """

    title: str = None
    content: str = None
    slide: Slide = None
    theme: str = None
    font_name: str = "Times New Roman"
    font_title_color: rgb = RGBColor(0, 0, 0)
    font_title_size: Pt = Pt(40)
    title_box = (Cm(2.81), Cm(5.44), Cm(21.59), Cm(4.08))

    def __init__(self, presentation, theme_path, title, *args, **kwargs):
        self.presentation = presentation
        self.slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.title = title
        self.theme = theme_path
        self.theme_param_path = os.path.join(self.theme, "mode.json")
        with open(self.theme_param_path, encoding="utf-8") as f:
            self.theme_param = json.load(f)
        print(self.theme)
        first_page_params = self.theme_param["first_page"]
        self.title_box = (
            Cm(first_page_params["title_info"]["pos_x"]),
            Cm(first_page_params["title_info"]["pos_y"]),
            Cm(first_page_params["title_info"]["width"]),
            Cm(first_page_params["title_info"]["height"]),
        )
        self.font_title_size = Pt(first_page_params["title_info"]["font_size"])
        self.font_name = first_page_params["title_info"]["font_name"]
        self.img_theme = theme_path + "/" + "img"
        self.init_font(**kwargs)
        self.init_slide()
        self.init_title()

    def init_slide(self):
        if os.path.exists(os.path.join(self.theme, "title.jpg")):
            path = os.path.join(self.theme, "title.jpg")
        elif os.path.exists(os.path.join(self.theme, "title.png")):
            path = os.path.join(self.theme, "title.png")
        else:
            path = get_random_file(self.img_theme)
        left, top, width, height = (Cm(0), Cm(0), Cm(25.4), Cm(14.29))
        picture = self.slide.shapes.add_picture(path, left, top, width, height)
        # set the width and height
        picture.left = 0
        picture.top = 0
        picture.width = self.presentation.slide_width
        picture.height = self.presentation.slide_height

    def init_font(self, **kwargs):
        if "font_name" in kwargs:
            self.font_name = kwargs["font_name"]
        if "font_title_size" in kwargs:
            self.font_title_size = kwargs["font_title_size"]
        if "font_content_size" in kwargs:
            self.font_content_size = kwargs["font_content_size"]
        if "font_title_color" in kwargs:
            self.font_title_color = kwargs["font_title_color"]

    def get_font(self, font: Font, category: str):
        font.name = self.font_name
        if category == MarkdownCategory.TITLE:
            font.size = self.font_title_size
            font.color.rgb = self.font_title_color
            font.bold = True

    def init_title(self):
        shapes = self.slide.shapes
        text_box = shapes.add_textbox(*self.title_box)
        tf = text_box.text_frame
        tf.clear()  # Clear existing content
        # tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        # tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True  # allow text to wrap

        # add title
        paragraph = tf.paragraphs[0]
        # paragraph.alignment = PP_ALIGN.CENTER
        paragraph.text = self.title
        self.get_font(paragraph.font, MarkdownCategory.TITLE)
        paragraph.word_wrap = True
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
