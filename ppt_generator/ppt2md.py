"""
This module extracts text and images from a PPTX file and converts them into a Markdown file with image placeholders.
It uses slide layout position to determine whether a slide is a Theme slide or a Paper slide.
"""

import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PptToMarkdownConverter:
    def __init__(self, ppt_path, theme_path, output_img_dir=None, output_md_path=None):
        """
        Args:
            ppt_path (str): Input PPTX file path.
            output_md_path (str): Output Markdown file path.
            output_img_dir (str or None): Output image directory. If None, do not extract images.
            theme_path (str): Parent path to the mode.json that defines theme page title_box.
        """
        self.ppt_path = ppt_path
        self.output_md_path = output_md_path
        self.output_img_dir = output_img_dir
        self.theme_mode_path = os.path.join(theme_path, "mode.json")

        self.slides_info = []  # Will store list of {'title': str, 'content': str, 'images': [img1, img2, ...], 'is_theme': bool}

        # Detect if images need to be saved
        self.save_images = output_img_dir is not None and output_img_dir != ""
        if self.save_images and not os.path.exists(self.output_img_dir):
            os.makedirs(self.output_img_dir)

        with open(self.theme_mode_path, "r", encoding="utf-8") as f:
            theme_param = json.load(f)
        theme_page = theme_param["theme_page"]
        self.title_box = (
            float(theme_page["title_info"]["pos_x"]),
            float(theme_page["title_info"]["pos_y"]),
            float(theme_page["title_info"]["width"]),
            float(theme_page["title_info"]["height"]),
        )

    def convert(self, isSave=False):
        self._extract_text_and_images()
        markdown_text = self._generate_markdown()
        if isSave: self._save_markdown(markdown_text) # only save when needed
        return markdown_text

    def _extract_text_and_images(self):
        presentation = Presentation(self.ppt_path)
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        for idx, slide in enumerate(presentation.slides):
            slide_info = {"title": "", "content": "", "images": [], "is_theme": False}
            img_count = 0

            for shape in slide.shapes:
                if self.save_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if self._is_background_image(shape, slide_width, slide_height):
                        continue  # Skip background images
                    img_count += 1
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    img_filename = f"slide{idx+1}_img{img_count}.{ext}"
                    img_path = os.path.join(self.output_img_dir, img_filename)
                    with open(img_path, 'wb') as f:
                        f.write(image_bytes)
                    # slide_info["images"].append(img_filename)
                    slide_info["images"].append(img_path)

                elif shape.has_text_frame:
                    text = shape.text.strip()
                    if not text:
                        continue
                    text = text.replace("\u000b", "\n").replace("\r", "\n")

                    # First non-empty text as title
                    if slide_info["title"] == "":
                        slide_info["title"] = text
                        # Record position to determine Theme slides
                        left_cm = shape.left.cm
                        top_cm = shape.top.cm
                        if self._is_within_title_box(left_cm, top_cm):
                            slide_info["is_theme"] = True
                    else:
                        if slide_info["content"]:
                            slide_info["content"] += "\n" + text
                            # slide_info["content"] += text
                        else:
                            slide_info["content"] = text

            self.slides_info.append(slide_info)

    def _is_within_title_box(self, left_cm, top_cm, tolerance_cm=1.0):
        """
        Determines if the given (left, top) falls within the theme page title box area.
        """
        title_left, title_top, title_width, title_height = self.title_box
        return (abs(left_cm - title_left) <= tolerance_cm and
                abs(top_cm - title_top) <= tolerance_cm)
    
    def _is_background_image(self, shape, slide_width, slide_height, tolerance=20000):
        """
        Determines if the given shape is a background image based on its position and size.
        
        Args:
            shape: pptx.shape
            slide_width: slide width (Emu)
            slide_height: slide height (Emu)
            tolerance: allowed deviation in EMU (default 20000 EMU ≈ 0.25 cm)
        """
        return (
            abs(shape.left) <= tolerance and
            abs(shape.top) <= tolerance and
            abs(shape.width - slide_width) <= tolerance and
            abs(shape.height - slide_height) <= tolerance
        )

    def _generate_markdown(self):
        md_lines = []

        for idx, slide in enumerate(self.slides_info):
            title = slide["title"].strip()
            content = slide["content"].strip()
            if self.save_images: images = slide["images"]
            is_theme = slide["is_theme"]

            if idx == 0:
                # First slide: main title
                md_lines.append(f"# {title}\n")
            elif title == "Table of Contents":
                # Skip Content page
                continue
            else:
                if is_theme:
                    # Theme slide
                    md_lines.append(f"## {title}\n")
                else:
                    # Paper slide
                    md_lines.append(f"### {title}\n")
                    if content:
                        md_lines.append(content)
                        
                    if self.save_images: 
                        for img_filename in images:
                            absolute_img_path = os.path.abspath(img_filename)
                            md_lines.append(f"![inserted_image]({absolute_img_path})")

            # md_lines.append("")  # Blank line between sections
            # only add blank line after paper (not theme page and title page)
            if not is_theme and idx != len(self.slides_info) - 1 and idx != 0:
                md_lines.append("")

        return "\n".join(md_lines)

    def _save_markdown(self, markdown_text):
        if not self.output_md_path:
            print("output_md_path not designated!")
            return
        with open(self.output_md_path, "w", encoding="utf-8") as f:
            f.write(markdown_text)
